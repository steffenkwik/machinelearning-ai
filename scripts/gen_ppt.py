# -*- coding: utf-8 -*-
"""Generate the StuntCare AI presentation deck (16:9), beginner-friendly."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK = RGBColor(0x05, 0x31, 0x2A)
GREEN = RGBColor(0x03, 0x98, 0x55)
MINT = RGBColor(0x6C, 0xE9, 0xA6)
LIGHT = RGBColor(0xEA, 0xF7, 0xF0)
INK = RGBColor(0x15, 0x24, 0x1C)
GREY = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xB4, 0x54, 0x09)
AMBERBG = RGBColor(0xFE, 0xF3, 0xC7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rect(sl, x, y, w, h, color):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    tb = sl.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt; f = r.font
            f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic
            f.name = "Calibri"
    return tb


def bullets(sl, x, y, w, h, items, size=18, color=INK, gap=10):
    tb = sl.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        rb = p.add_run(); rb.text = "▸  "; rb.font.size = Pt(size); rb.font.color.rgb = GREEN; rb.font.bold = True
        for j, part in enumerate(it.split("**")):
            if not part:
                continue
            r = p.add_run(); r.text = part; r.font.size = Pt(size)
            r.font.color.rgb = color; r.font.bold = (j % 2 == 1); r.font.name = "Calibri"
    return tb


def infobox(sl, x, y, w, h, title, body, bg=LIGHT, accent=GREEN, tcolor=None):
    rect(sl, x, y, w, h, bg)
    rect(sl, x, y, Inches(0.09), h, accent)
    runs = [[(title, 13, accent, True, False)]]
    for line in body.split("\n"):
        runs.append([(line, 11.5, tcolor or INK, False, False)])
    text(sl, x + Inches(0.22), y + Inches(0.12), w - Inches(0.4), h - Inches(0.22), runs, space=3)


def content_slide(title, eyebrow=None):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, SW, SH, WHITE)
    rect(sl, 0, 0, Inches(0.22), SH, GREEN)
    rect(sl, Inches(0.6), Inches(0.5), Inches(0.6), Inches(0.09), MINT)
    if eyebrow:
        text(sl, Inches(0.6), Inches(0.6), Inches(11.5), Inches(0.4),
             [[(eyebrow.upper(), 12, GREEN, True, False)]])
    text(sl, Inches(0.6), Inches(0.88), Inches(12.2), Inches(1.0),
         [[(title, 28, DARK, True, False)]])
    return sl


def stat_card(sl, x, y, w, value, label):
    rect(sl, x, y, w, Inches(1.6), LIGHT)
    rect(sl, x, y, w, Inches(0.09), GREEN)
    text(sl, x, y + Inches(0.26), w, Inches(0.8), [[(value, 32, GREEN, True, False)]], align=PP_ALIGN.CENTER)
    text(sl, x, y + Inches(1.0), w, Inches(0.5), [[(label, 12, GREY, False, False)]], align=PP_ALIGN.CENTER)


# ───────────── Slide 1: Title ─────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, DARK)
rect(sl, 0, 0, SW, Inches(0.18), GREEN)
rect(sl, 0, SH - Inches(0.18), SW, Inches(0.18), GREEN)
rect(sl, Inches(0.9), Inches(1.5), Inches(0.7), Inches(0.11), MINT)
text(sl, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     [[("SDG #3 · GOOD HEALTH AND WELL-BEING", 14, MINT, True, False)]])
text(sl, Inches(0.9), Inches(2.3), Inches(11.6), Inches(2.2), [
    [("StuntCare AI", 54, WHITE, True, False)],
    [("Deteksi Dini Risiko Stunting Balita Multi-Faktor", 26, LIGHT, False, False)],
    [("berbasis Random Forest yang berjalan di peramban", 26, LIGHT, False, False)],
])
text(sl, Inches(0.9), Inches(5.5), Inches(11.5), Inches(1.3), [
    [("Daniel Steffen K  ·  NIM 2602071171", 18, WHITE, True, False)],
    [("Artificial Intelligence (COMP6065001) · LA05-LEC · BINUS University", 14, MINT, False, False)],
])

# ───────────── Slide 2: Masalah (+ KEK→BBLR) ─────────────
sl = content_slide("Stunting: gagal tumbuh karena banyak faktor", "Latar belakang")
bullets(sl, Inches(0.65), Inches(1.95), Inches(7.3), Inches(3.2), [
    "Prevalensi stunting Indonesia **19,8%** (SSGI 2024) — ±4,4 juta balita.",
    "Stunting = anak **lebih pendek** dari standar usianya akibat kurang gizi kronis.",
    "Dampak jangka panjang: kecerdasan, daya tahan tubuh, dan produktivitas.",
    "Penyebabnya **banyak faktor** — bukan hanya tinggi badan anak.",
], size=17)
infobox(sl, Inches(0.65), Inches(5.25), Inches(11.9), Inches(1.55),
        "Apa itu \"KEK → BBLR\"?",
        "KEK = Kurang Energi Kronis: ibu kekurangan gizi dalam waktu lama saat hamil. "
        "Ini bisa menyebabkan BBLR = Berat Badan Lahir Rendah (bayi lahir < 2,5 kg). "
        "Bayi BBLR lebih berisiko mengalami stunting. Karena itu gizi ibu saat hamil ikut dihitung.")
stat_card(sl, Inches(8.3), Inches(1.95), Inches(4.25), "19,8%", "Prevalensi nasional (SSGI 2024)")
stat_card(sl, Inches(8.3), Inches(3.7), Inches(4.25), "7", "Faktor penyebab dianalisis")

# ───────────── Slide 3: Solusi (+ peramban + z-score) ─────────────
sl = content_slide("Solusi: prediksi 7 faktor, langsung di perangkat", "Gagasan")
bullets(sl, Inches(0.65), Inches(1.9), Inches(12), Inches(2.0), [
    "Aplikasi web yang menggabungkan **7 faktor** menjadi satu prediksi.",
    "Mesin prediksi: **Random Forest** (akurasi 88,40%).",
    "Output: status gizi + tingkat keyakinan + rekomendasi tindak lanjut.",
], size=17)
infobox(sl, Inches(0.65), Inches(4.05), Inches(5.85), Inches(2.75),
        "Apa itu \"berjalan 100% di peramban\"?",
        "Peramban = aplikasi browser (Chrome, Safari, dll). Artinya semua perhitungan AI "
        "terjadi di HP/laptop pengguna sendiri — BUKAN di komputer server kami. "
        "Keuntungan: gratis, cepat, bisa tanpa internet, dan data anak tidak dikirim ke mana pun.")
infobox(sl, Inches(6.7), Inches(4.05), Inches(5.85), Inches(2.75),
        "Apa itu \"Z-Score WHO (HAZ)\"?",
        "Angka yang membandingkan tinggi anak dengan tinggi standar WHO untuk anak seumur & "
        "sejenis kelamin. 0 = persis rata-rata; di bawah −2 SD = lebih pendek dari normal "
        "(indikasi stunting); di atas +2 = lebih tinggi dari rata-rata.")

# ───────────── Slide 4: 7 Faktor (+ BKKBN 4T) ─────────────
sl = content_slide("Tujuh faktor masukan model", "Fitur")
left = ["**Tinggi badan** (cm) — indikator utama",
        "**Umur** (bulan) — acuan standar WHO",
        "**Berat badan** (kg)",
        "**Jenis kelamin** — standar L/P beda"]
right = ["**Jarak kehamilan** (bulan)",
         "**Usia ibu menikah** (tahun)",
         "**Gizi ibu saat hamil** (Baik/Sedang/Buruk)"]
bullets(sl, Inches(0.65), Inches(1.95), Inches(5.9), Inches(3.0), left, size=16)
bullets(sl, Inches(6.7), Inches(1.95), Inches(5.9), Inches(3.0), right, size=16)
infobox(sl, Inches(0.65), Inches(5.25), Inches(11.9), Inches(1.55),
        "Apa itu \"ideal > 24 bulan (BKKBN 4T)\"?",
        "Program \"4 Terlalu\" BKKBN: Terlalu muda, Terlalu tua, Terlalu dekat, Terlalu banyak. "
        "Jarak kehamilan ideal > 24 bulan agar tubuh ibu pulih dulu. Jarak terlalu dekat "
        "menaikkan risiko stunting pada anak berikutnya.")

# ───────────── Slide 5: Random Forest cara kerja ─────────────
sl = content_slide("Random Forest — bagaimana ia bekerja", "Inti algoritma")
bullets(sl, Inches(0.65), Inches(1.9), Inches(7.5), Inches(4.6), [
    "**Ensemble** = gabungan **120 pohon keputusan** (decision tree).",
    "Tiap pohon belajar dari **sebagian data acak** (teknik bagging).",
    "Tiap pohon memberi \"suara\" prediksi; hasil akhir = **rata-rata** semua pohon.",
    "Banyak pohon → hasil **lebih akurat & stabil** (tidak mudah salah/overfit).",
    "Bonus: bisa menunjukkan **faktor mana yang paling penting** (transparan).",
], size=16)
for i in range(3):
    x = Inches(8.7 + i * 1.4)
    rect(sl, x, Inches(2.5), Inches(1.1), Inches(2.5), LIGHT)
    rect(sl, x, Inches(2.5), Inches(1.1), Inches(0.08), GREEN)
    text(sl, x, Inches(3.35), Inches(1.1), Inches(0.6),
         [[(f"Pohon {i+1}", 14, GREEN, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(sl, Inches(8.7), Inches(5.2), Inches(4.0), Inches(0.6),
     [[("120 pohon → voting → 1 prediksi", 13, GREY, False, True)]], align=PP_ALIGN.CENTER)

# ───────────── Slide 6: Pelatihan + encoding (beginner) ─────────────
sl = content_slide("Bagaimana model dilatih — langkah demi langkah", "Inti algoritma")
bullets(sl, Inches(0.65), Inches(1.85), Inches(7.3), Inches(4.8), [
    "**1. Kumpulkan data:** 25.000 contoh (7 faktor + status gizinya).",
    "**2. Encoding:** ubah teks menjadi angka agar bisa dihitung komputer.",
    "**3. Bagi data:** 80% untuk belajar (latih), 20% untuk ujian (uji).",
    "**4. Latih:** 120 pohon belajar pola dari data latih.",
    "**5. Uji:** model diuji pada 20% data yang belum pernah dilihat.",
], size=16)
infobox(sl, Inches(8.2), Inches(1.85), Inches(4.4), Inches(2.45),
        "Apa itu \"Encoding\"?",
        "Komputer hanya mengerti angka, bukan teks. Encoding = mengubah teks jadi angka. "
        "Contoh: Laki-laki→0, Perempuan→1; gizi Baik→0, Buruk→1, Sedang→2. "
        "Ini disebut Label Encoding.")
box = rect(sl, Inches(8.2), Inches(4.5), Inches(4.4), Inches(2.0), DARK)
text(sl, Inches(8.4), Inches(4.66), Inches(4.05), Inches(1.7), [
    [("rf = RandomForestClassifier(", 12, MINT, False, False)],
    [("      n_estimators=120,", 12, WHITE, False, False)],
    [("      max_depth=14)", 12, WHITE, False, False)],
    [("rf.fit(X_train, y_train)", 12, WHITE, False, False)],
], space=2)
text(sl, Inches(8.2), Inches(6.55), Inches(4.4), Inches(0.4),
     [[("Kode inti pelatihan (scikit-learn)", 11, GREY, False, True)]])

# ───────────── Slide 7: Pengujian / Performa (beginner) ─────────────
sl = content_slide("Bagaimana kami menguji akurasinya", "Hasil")
bullets(sl, Inches(0.65), Inches(1.9), Inches(12), Inches(2.2), [
    "Model diuji pada **5.000 data ujian** yang TIDAK dipakai saat latihan — "
    "seperti soal ujian baru, supaya penilaian jujur.",
    "**Akurasi 88,40%** artinya: dari 100 anak, sekitar **88 diprediksi benar**.",
    "**Cross-validation:** uji berulang 5 kali dengan pembagian berbeda; hasil stabil "
    "(±0,42%) → model benar-benar belajar pola, bukan sekadar menghafal.",
], size=15)
stat_card(sl, Inches(0.65), Inches(4.55), Inches(3.85), "88,40%", "Akurasi (Excellent ≥85%)")
stat_card(sl, Inches(4.75), Inches(4.55), Inches(3.85), "83,56%", "F1-Score (lihat catatan)")
stat_card(sl, Inches(8.85), Inches(4.55), Inches(3.7), "87,31%", "Cross-Validation (5x)")
text(sl, Inches(0.65), Inches(6.35), Inches(12), Inches(0.6),
     [[("F1-Score = ukuran yang adil saat jumlah tiap kelas tidak seimbang (kelas Normal jauh lebih banyak).",
        13, GREY, False, True)]])

# ───────────── Slide 8: Feature importance + siapa yang menentukan ─────────────
sl = content_slide("Faktor paling berpengaruh — siapa yang menentukan?", "Hasil")
fi = [("Tinggi badan", 50.6), ("Umur", 28.7), ("Berat badan", 9.5),
      ("Jarak kehamilan", 4.3), ("Usia ibu menikah", 4.2), ("Gizi ibu hamil", 1.5), ("Jenis kelamin", 1.3)]
y = Inches(1.95); maxw = 5.2
for name, val in fi:
    text(sl, Inches(0.65), y, Inches(2.6), Inches(0.34), [[(name, 13, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(sl, Inches(3.35), y + Inches(0.04), Inches(maxw * val / 50.6), Inches(0.26), GREEN)
    text(sl, Inches(3.35 + maxw * val / 50.6 + 0.1), y, Inches(1.0), Inches(0.34),
         [[(f"{val}%", 12, GREY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)
infobox(sl, Inches(9.0), Inches(1.95), Inches(3.55), Inches(4.05),
        "Siapa yang menentukan?",
        "BUKAN kami. Random Forest MENGHITUNG SENDIRI dari data: seberapa sering & seberapa "
        "besar tiap faktor membantu menebak status gizi dengan benar. Hasilnya disebut "
        "\"feature importance\". Tinggi & umur paling besar karena definisi stunting memang "
        "soal tinggi-menurut-umur.")

# ───────────── Slide 9: Teknologi yang kami bangun (no Streamlit) ─────────────
sl = content_slide("Teknologi yang kami bangun", "Rekayasa perangkat lunak")
bullets(sl, Inches(0.65), Inches(1.95), Inches(11.9), Inches(4.6), [
    "**React + Vite + Tailwind** — kerangka modern untuk membuat tampilan website "
    "yang cepat dan rapi.",
    "**ONNX Runtime Web** — pustaka yang bisa menjalankan model machine learning langsung "
    "di dalam peramban (browser), tanpa server.",
    "**Web Worker** — \"pekerja latar belakang\"; proses menebak dijalankan terpisah agar "
    "tampilan tidak macet/hang saat menghitung.",
    "**Hasil:** aplikasi gratis, cepat, bisa offline, dan data anak tidak keluar dari perangkat.",
], size=17)

# ───────────── Slide 10: Arsitektur (beginner) ─────────────
sl = content_slide("Arsitektur: cara semua bekerja bersama", "Arsitektur")
steps = [("1. Isi 7 faktor", "Pengguna mengisi data anak & ibu di form."),
         ("2. Web Worker", "Data diproses di pekerja latar, di perangkat."),
         ("3. Random Forest", "Model (format ONNX) menghitung dalam sekejap."),
         ("4. Hasil", "Status gizi + keyakinan + rekomendasi tampil.")]
x = Inches(0.65)
for t, d in steps:
    rect(sl, x, Inches(2.0), Inches(2.85), Inches(2.0), LIGHT)
    rect(sl, x, Inches(2.0), Inches(2.85), Inches(0.09), GREEN)
    text(sl, x + Inches(0.15), Inches(2.25), Inches(2.55), Inches(0.6), [[(t, 15, DARK, True, False)]])
    text(sl, x + Inches(0.15), Inches(2.85), Inches(2.55), Inches(1.0), [[(d, 12, GREY, False, False)]])
    x += Inches(3.02)
infobox(sl, Inches(0.65), Inches(4.45), Inches(11.9), Inches(1.5),
        "Istilah teknis (sederhana)",
        "ONNX = format standar agar model AI bisa dipakai di mana saja. "
        "WASM (WebAssembly) = teknologi yang membuat kode berjalan cepat di dalam peramban. "
        "Gabungan keduanya membuat model berjalan instan, di perangkat, dan bisa offline.")
text(sl, Inches(0.65), Inches(6.15), Inches(12), Inches(0.5),
     [[("Tanpa server, tanpa biaya, dan privasi terjaga (data tidak dikirim ke mana pun).",
        13, GREY, False, True)]])

# ───────────── Slide 11: Fitur foto (beginner) ─────────────
sl = content_slide("Fitur opsional: Analisis Foto (eksperimental)", "Pengayaan")
bullets(sl, Inches(0.65), Inches(1.9), Inches(7.4), Inches(4.7), [
    "Pengguna bisa **unggah foto** atau pakai **kamera** anak.",
    "**MediaPipe** (pustaka Google) mendeteksi 33 titik tubuh + siluet badan, "
    "langsung di perangkat (foto tidak diunggah ke server).",
    "Dihitung **proporsi tubuh** (kurus/normal/gemuk) dengan mengukur ketebalan "
    "lengan/kaki dibanding panjangnya — bekerja saat berdiri, duduk, atau berbaring.",
], size=15)
infobox(sl, Inches(8.3), Inches(1.9), Inches(4.25), Inches(4.55),
        "Kenapa hanya \"pendukung\"?",
        "JUJUR: foto 2D tidak bisa mengukur tinggi badan, jadi BUKAN penentu stunting. "
        "Hasil foto hanya menggeser sedikit probabilitas akhir (maksimum ±10 poin) dan "
        "ditampilkan transparan (sebelum → sesudah). Random Forest TIDAK dilatih ulang — "
        "jadi tidak ada klaim palsu. Berguna untuk mendeteksi kurus/gemuk.",
        bg=AMBERBG, accent=AMBER, tcolor=RGBColor(0x6B, 0x4A, 0x07))

# ───────────── Slide 12: Fitur chat AI (NEW) ─────────────
sl = content_slide("Fitur baru: Asisten AI Gizi (chat)", "Inovasi tambahan")
bullets(sl, Inches(0.65), Inches(1.9), Inches(7.4), Inches(4.7), [
    "Di **atas form**, pengguna bisa bertanya ke asisten AI tentang **cara mencegah "
    "stunting & tips gizi** — agar dapat ilmu, bukan hanya hasil prediksi.",
    "Ditenagai **Claude (Anthropic)**. Jawaban ringkas, berbasis pedoman WHO/Kemenkes, "
    "dan menolak memberi diagnosis medis.",
    "Pertanyaan dikirim ke layanan AI; **prediksi stunting tetap 100% di perangkat**.",
], size=15)
infobox(sl, Inches(8.3), Inches(1.9), Inches(4.25), Inches(4.7),
        "Langkah integrasi (singkat)",
        "1. Buat API key di console.anthropic.com.\n"
        "2. Simpan sebagai secret ANTHROPIC_API_KEY di Cloudflare Pages.\n"
        "3. Permintaan chat lewat fungsi server (Cloudflare Pages Function) agar API key "
        "aman — tidak pernah ada di sisi pengguna.\n"
        "4. Deploy ulang. Selesai.")

# ───────────── Slide 13: Demo & deploy ─────────────
sl = content_slide("Demo & deployment", "Implementasi")
bullets(sl, Inches(0.65), Inches(2.0), Inches(11.9), Inches(3.0), [
    "Live di **Cloudflare Pages** (statis, gratis, HTTPS otomatis).",
    "URL: **machinelearning-ai.pages.dev**",
    "Build: Root = web · Build command = npm run build · Output = dist.",
    "Repositori: github.com/steffenkwik/machinelearning-ai",
], size=17)
rect(sl, Inches(0.65), Inches(5.2), Inches(11.9), Inches(1.1), LIGHT)
text(sl, Inches(0.85), Inches(5.4), Inches(11.5), Inches(0.8),
     [[("Demo langsung: isi 7 faktor → tekan Analisis → hasil tampil seketika (offline).",
        16, DARK, True, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ───────────── Slide 14: Kesimpulan ─────────────
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, SW, SH, DARK)
rect(sl, 0, 0, SW, Inches(0.18), GREEN)
text(sl, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.6), [[("KESIMPULAN", 16, MINT, True, False)]])
bd = [
    "**Random Forest 88,40%** menggabungkan 7 faktor penyebab stunting.",
    "Model berjalan **100% di peramban** — gratis, cepat, offline, dan privat.",
    "Dilengkapi **analisis foto** & **asisten AI gizi** sebagai pengayaan yang jujur.",
    "Layak sebagai **alat skrining awal**; lanjutan: validasi dengan data klinis nyata.",
]
tb = sl.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(3.6))
tf = tb.text_frame; tf.word_wrap = True
for i, it in enumerate(bd):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(14)
    rb = p.add_run(); rb.text = "▸  "; rb.font.size = Pt(19); rb.font.color.rgb = MINT; rb.font.bold = True
    for j, part in enumerate(it.split("**")):
        if not part:
            continue
        r = p.add_run(); r.text = part; r.font.size = Pt(19); r.font.color.rgb = WHITE; r.font.bold = (j % 2 == 1)
text(sl, Inches(0.9), Inches(5.85), Inches(11.5), Inches(1.0), [
    [("Terima kasih.", 30, WHITE, True, False)],
    [("Daniel Steffen K · NIM 2602071171 · BINUS University", 14, MINT, False, False)],
])

base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
out = os.path.join(base, "StuntCare_AI_Presentation.pptx")
prs.save(out)
print("Slides:", len(prs.slides._sldIdLst), "->", out)
